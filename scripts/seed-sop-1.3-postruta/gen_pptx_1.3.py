#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Genera la PPT de capacitación del SOP 1.3 — Procesos post-ruta y cierre físico.
Clona el lenguaje visual de la PPT 1.2 (navy MERCOSUR·DPO, ámbar, Calibri, 16:9).
Salida: /tmp/material-1.3.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Paleta (idéntica a la 1.2) ---
NAVY   = RGBColor(0x0A, 0x16, 0x28)   # fondo
CARD   = RGBColor(0x1E, 0x2A, 0x3C)   # tarjetas
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # acento
SLATE  = RGBColor(0x94, 0xA3, 0xB8)   # texto secundario
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0xB9, 0x81)   # OK
RED    = RGBColor(0xEF, 0x44, 0x44)   # NO
FONT   = "Calibri"

EMU = 914400
SW, SH = 13.333, 7.5  # 16:9
TOTAL = 15
FOOTER = "SOP 1.3 — Procesos post-ruta y cierre físico  ·  Pilar Entrega  ·  Mercosur DPO"

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]


def _in(v): return Emu(int(v * EMU))


def slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        _in(x), _in(y), _in(w), _in(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: lista de párrafos; cada párrafo = lista de (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            run = p.add_run(); run.text = txt
            run.font.name = FONT; run.font.size = Pt(size)
            run.font.color.rgb = color; run.font.bold = bold
    return tb


def chrome(s, title, subtitle, n):
    """Barra superior: título + subtítulo + número de slide + footer."""
    box(s, 0.7, 0.55, 0.12, 0.95, fill=AMBER)  # barra ámbar lateral
    text(s, 0.95, 0.5, 9.5, 1.1, [
        [(title, 30, WHITE, True)],
        [(subtitle, 15, SLATE, False)],
    ], space_after=2)
    text(s, 11.3, 0.6, 1.4, 0.4, [[(f"{n} / {TOTAL}", 12, SLATE, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.95, 7.02, 11.4, 0.35, [[(FOOTER, 9, SLATE, False)]])


def bullets(s, x, y, w, h, items, size=14, color=WHITE, dash="—  ", gap=6):
    runs = [[(dash, size, AMBER, True), (it, size, color, False)] for it in items]
    text(s, x, y, w, h, runs, space_after=gap, line_spacing=1.05)


# ============================ SLIDE 1 — Portada ============================
s = slide()
box(s, 0, 0, SW, 0.18, fill=AMBER)
text(s, 1.0, 2.0, 11, 0.5, [[("MERCOSUR  ·  DPO", 16, AMBER, True)]])
text(s, 1.0, 2.6, 11.3, 1.6, [
    [("SOP 1.3 — Procesos post-ruta", 40, WHITE, True)],
    [("y cierre físico", 40, WHITE, True)],
], space_after=2, line_spacing=1.0)
text(s, 1.0, 4.5, 11, 0.5, [[("Pilar Entrega  ·  Capacitación operativa para choferes y ayudantes", 16, SLATE, False)]])
box(s, 1.0, 5.35, 6.2, 0.02, fill=CARD)
text(s, 1.0, 5.55, 11, 0.5, [[("Instructor: Fausto Azzaretti  ·  Duración: 1 h", 14, SLATE, False)]])

# ===================== SLIDE 2 — Objetivo y alcance =====================
s = slide(); chrome(s, "Objetivo y alcance", "Por qué hacemos este SOP", 2)
box(s, 0.95, 1.85, 5.55, 4.7, fill=CARD, round_=True)
text(s, 1.3, 2.15, 4.9, 0.4, [[("OBJETIVO", 15, AMBER, True)]])
text(s, 1.3, 2.7, 4.9, 3.6, [[("Definir los pasos del proceso de retorno de camiones al CD: el cierre financiero, la verificación física, el control de envases y devoluciones, y la revisión de la unidad, garantizando seguridad, control y trazabilidad.", 15, WHITE, False)]], line_spacing=1.15)
box(s, 6.8, 1.85, 5.55, 4.7, fill=CARD, round_=True)
text(s, 7.15, 2.15, 4.9, 0.4, [[("ALCANCE", 15, AMBER, True)]])
text(s, 7.15, 2.7, 4.9, 3.6, [[("Aplica a todo el equipo de distribución: choferes, ayudantes, supervisor de rutas (SDR), autoelevadoristas, analistas y asistentes de distribución que intervienen en el retorno al CD.", 15, WHITE, False)]], line_spacing=1.15)

# ============================ SLIDE 3 — RACI ============================
s = slide(); chrome(s, "RACI", "Roles y responsabilidades", 3)
raci = [("R  —  Responsable", "El que ejecuta la acción"),
        ("A  —  Accountable", "Dueño del KPI, último responsable"),
        ("C  —  Consultado", "A quien se consulta antes de la decisión"),
        ("I  —  Informado", "A quien se informa después de la acción")]
for i, (h, d) in enumerate(raci):
    x = 0.95 + (i % 2) * 5.85
    y = 1.95 + (i // 2) * 2.35
    box(s, x, y, 5.55, 2.05, fill=CARD, round_=True)
    text(s, x + 0.35, y + 0.3, 5.0, 0.4, [[(h, 19, AMBER, True)]])
    text(s, x + 0.35, y + 1.0, 5.0, 0.8, [[(d, 14, WHITE, False)]], line_spacing=1.1)

# ============================ SLIDE 4 — EPP ============================
s = slide(); chrome(s, "Elementos de Protección Personal", "Obligatorios en toda la jornada", 4)
epp = [("Guantes anticortes", "Manipulación de cajones y envases"),
       ("Botines de seguridad", "Punta de acero, suela antideslizante"),
       ("Gafas de seguridad", "Descarga y manipulación"),
       ("Faja lumbar", "Soporte ergonómico al levantar peso"),
       ("Uniforme de trabajo", "Vestimenta limpia y reflectiva"),
       ("Ropa reflectiva", "Visibilidad en el CD y la maniobra")]
for i, (h, d) in enumerate(epp):
    x = 0.95 + (i % 3) * 3.92
    y = 2.0 + (i // 3) * 2.25
    box(s, x, y, 3.65, 1.95, fill=CARD, round_=True)
    text(s, x + 0.3, y + 0.28, 3.1, 0.4, [[(h, 15, AMBER, True)]])
    text(s, x + 0.3, y + 0.95, 3.1, 0.85, [[(d, 12.5, SLATE, False)]], line_spacing=1.05)

# ===================== SLIDE 5 — Definiciones clave =====================
s = slide(); chrome(s, "Definiciones clave", "Vocabulario del proceso", 5)
defs = [("Equipo de Distribución", "Choferes y ayudantes que ejecutan la ruta"),
        ("SDR", "Supervisor de Rutas"),
        ("Recepcionista", "Recibe los camiones al regreso al CD"),
        ("Autoelevadorista", "Conduce el autoelevador en la descarga"),
        ("Sector de vacíos", "Zona definida para la descarga de camiones"),
        ("Rechazo", "Producto que el cliente no recibió")]
for i, (h, d) in enumerate(defs):
    x = 0.95 + (i % 2) * 5.85
    y = 1.9 + (i // 2) * 1.62
    box(s, x, y, 5.55, 1.4, fill=CARD, round_=True)
    text(s, x + 0.32, y + 0.22, 5.0, 0.4, [[(h, 14.5, AMBER, True)]])
    text(s, x + 0.32, y + 0.74, 5.0, 0.5, [[(d, 12.5, WHITE, False)]])

# ============== SLIDE 6 — El retorno: dos caminos paralelos ==============
s = slide(); chrome(s, "El retorno al CD", "Dos procesos en paralelo", 6)
box(s, 0.95, 2.0, 5.55, 4.4, fill=CARD, round_=True)
text(s, 1.3, 2.35, 4.9, 0.4, [[("CHOFER", 13, SLATE, True)]])
text(s, 1.3, 2.8, 4.9, 0.6, [[("Cierre financiero", 24, AMBER, True)]])
bullets(s, 1.3, 3.75, 4.9, 2.5, [
    "Rinde el dinero de la jornada", "Controla NC, cheques y transferencias",
    "Completa la planilla de caja", "Entrega documentación a tesorería"], size=13.5, gap=7)
box(s, 6.8, 2.0, 5.55, 4.4, fill=CARD, round_=True)
text(s, 7.15, 2.35, 4.9, 0.4, [[("AYUDANTE", 13, SLATE, True)]])
text(s, 7.15, 2.8, 4.9, 0.6, [[("Verificación física", 24, AMBER, True)]])
bullets(s, 7.15, 3.75, 4.9, 2.5, [
    "Descarga de vacíos y devoluciones", "Control de envases y producto",
    "Limpieza y sanitizado de la unidad", "Checklist de retorno (CloudFleet)"], size=13.5, gap=7)

# ===================== SLIDE 7 — Cierre financiero =====================
s = slide(); chrome(s, "Cierre financiero", "Responsable: el chofer  ·  máx. 40 minutos", 7)
bullets(s, 0.95, 1.95, 11.5, 5.0, [
    "Solicita la llave de la caja fuerte al controlador, recolecta el dinero y devuelve la llave.",
    "Realiza el recuento del dinero con el contador de billetes y lo rinde a tesorería.",
    "Tesorería corrobora que el importe coincida con lo informado por el chofer.",
    "Controla las notas de crédito, cheques y transferencias recibidas.",
    "Completa la planilla de caja: efectivo, cheques, transferencias, cuentas corrientes, rechazos, cobranzas y venta de envases.",
    "Entrega la planilla y toda la documentación a tesorería y administración para la liquidación del reparto.",
], size=15, gap=11)

# ===================== SLIDE 8 — Verificación física =====================
s = slide(); chrome(s, "Verificación física", "Responsable: el ayudante  ·  máx. 40 minutos", 8)
bullets(s, 0.95, 1.85, 11.5, 5.1, [
    "Al ingresar al CD se toma la medición del odómetro y el equipo que retorna.",
    "Estaciona en el sector de descarga, coloca los tacos de seguridad y abre las cortinas del camión.",
    "Se ubica en la zona segura mientras los autoelevadores retiran los vacíos y devoluciones (supervisa permanentemente).",
    "Una vez finalizada la descarga, realiza el control de envases y devoluciones.",
    "Retira la unidad del parque de vacíos, la estaciona y realiza la limpieza y sanitizado (obligatorio diario).",
    "Completa el checklist de retorno (CloudFleet) y notifica combustible/calibración si corresponde.",
    "Deposita la llave en el buzón y marca el biométrico antes de salir del CD.",
], size=14.5, gap=8)

# ===================== SLIDE 9 — Control de envases =====================
s = slide(); chrome(s, "Control de envases", "Check-in y clasificación", 9)
cols = [("CHECK-IN (Chess)", ["El controlador ingresa por sistema los productos que retornan de la ruta: rechazos y envases."]),
        ("CLASIFICACIÓN", ["Se controlan envases y esqueletos (cajones).", "Se verifica la integridad de cada uno."]),
        ("COMPETENCIA", ["Cajones y envases de la competencia van a un sector exclusivo, separado de nuestra marca.", "Se evalúa con cervecería su canje periódico."])]
for i, (h, items) in enumerate(cols):
    x = 0.95 + i * 3.92
    box(s, x, 1.95, 3.65, 4.4, fill=CARD, round_=True)
    text(s, x + 0.3, 2.25, 3.1, 0.4, [[(h, 14, AMBER, True)]])
    bullets(s, x + 0.3, 2.95, 3.1, 3.2, items, size=12.5, gap=8)

# ===================== SLIDE 10 — Control de producto =====================
s = slide(); chrome(s, "Control de producto", "Rechazos y devoluciones", 10)
items = [("RECHAZO", "Producto que no pudo entregarse al PDV. Se registra y contabiliza para su seguimiento y análisis.", AMBER),
         ("DEVOLUCIÓN EN BUEN ESTADO", "Error de envío u otro motivo no ligado a calidad. Si no está próximo a vencer y no tiene daños → vuelve al stock.", GREEN),
         ("DEVOLUCIÓN EN MAL ESTADO", "Problema de packaging o frescura. Empaque roto → reempaque y vuelve al stock. Pinchados, rotos o sin gas → se desechan y descuentan.", RED)]
for i, (h, d, c) in enumerate(items):
    y = 1.9 + i * 1.62
    box(s, 0.95, y, 11.4, 1.42, fill=CARD, round_=True)
    box(s, 0.95, y, 0.1, 1.42, fill=c)
    text(s, 1.3, y + 0.18, 11.0, 0.4, [[(h, 14, c, True)]])
    text(s, 1.3, y + 0.66, 10.8, 0.7, [[(d, 13, WHITE, False)]], line_spacing=1.05)

# ===================== SLIDE 11 — Revisión de la unidad =====================
s = slide(); chrome(s, "Revisión de la unidad", "Diaria y mensual", 11)
box(s, 0.95, 1.95, 5.55, 4.4, fill=CARD, round_=True)
text(s, 1.3, 2.3, 4.9, 0.4, [[("DIARIA", 16, AMBER, True)]])
bullets(s, 1.3, 2.95, 4.9, 3.2, [
    "El chofer detecta problemas y los registra en CloudFleet (Checklist de retorno).",
    "El SDR coordina la resolución con el servicio técnico."], size=13.5, gap=9)
box(s, 6.8, 1.95, 5.55, 4.4, fill=CARD, round_=True)
text(s, 7.15, 2.3, 4.9, 0.4, [[("MENSUAL", 16, AMBER, True)]])
bullets(s, 7.15, 2.95, 4.9, 3.2, [
    "Una vez al mes el SDR completa un checklist digital.",
    "Revisa estado general y vencimientos (extintor, documentación).",
    "Carga la información en Smart 5S para programar actualizaciones."], size=13.5, gap=9)

# =============== SLIDE 12 — Limpieza del camión ===============
s = slide(); chrome(s, "Limpieza del camión", "Inocuidad de los alimentos transportados", 12)
box(s, 0.95, 2.1, 11.4, 3.4, fill=CARD, round_=True)
bullets(s, 1.4, 2.55, 10.6, 2.6, [
    "La carrocería debe estar siempre limpia para proteger los alimentos de la contaminación.",
    "Los productos de limpieza y desinfección tienen que ser aptos para uso alimentario.",
    "El agua utilizada debe ser potable.",
    "El sanitizado de carrocería e interior es obligatorio todos los días."], size=15, gap=12)

# =============== SLIDE 13 — KPIs / Tiempos internos ===============
s = slide(); chrome(s, "KPIs · Tiempos internos", "Cómo medimos el post-ruta", 13)
text(s, 0.95, 1.85, 11.4, 0.6, [[("Tiempo interno: desde que el equipo regresa de la ruta hasta que sale del CD.", 14, SLATE, False)]])
kpi = [("Verificación física", "Ayudante · descarga, control y limpieza", "≤ 40 min"),
       ("Cierre financiero", "Chofer · liquidación y rendición", "≤ 40 min")]
for i, (h, d, v) in enumerate(kpi):
    x = 0.95 + i * 5.85
    box(s, x, 2.7, 5.55, 3.4, fill=CARD, round_=True)
    text(s, x + 0.4, 3.05, 4.7, 0.4, [[(h, 18, AMBER, True)]])
    text(s, x + 0.4, 3.7, 4.7, 0.6, [[(d, 13, SLATE, False)]], line_spacing=1.05)
    text(s, x + 0.4, 4.6, 4.7, 1.1, [[(v, 44, WHITE, True)]])

# =============== SLIDE 14 — Rendimiento y feedback ===============
s = slide(); chrome(s, "Rendimiento y feedback de choferes", "Mejora continua del CD", 14)
box(s, 0.95, 2.0, 5.55, 4.3, fill=CARD, round_=True)
text(s, 1.3, 2.3, 4.9, 0.4, [[("RENDIMIENTO", 15, AMBER, True)]])
bullets(s, 1.3, 2.9, 4.9, 3.2, [
    "Gráficos diarios de performance de cada tripulación.",
    "Resumen del desempeño del día anterior, accesible desde el celular."], size=13.5, gap=9)
box(s, 6.8, 2.0, 5.55, 4.3, fill=CARD, round_=True)
text(s, 7.15, 2.3, 4.9, 0.4, [[("FEEDBACK", 15, AMBER, True)]])
bullets(s, 7.15, 2.9, 4.9, 3.2, [
    "Canal digital para reportar dificultades en la entrega (estacionamiento, demoras, ventana horaria, calles, cobro, maltrato).",
    "Seguimiento diario, planes de acción y comunicación en las matinales."], size=13.5, gap=9)

# ============================ SLIDE 15 — Cierre ============================
s = slide()
box(s, 0, 0, SW, 0.18, fill=AMBER)
text(s, 1.0, 1.3, 11.3, 0.6, [[("Cuestiones que recordar siempre", 30, WHITE, True)]])
bullets(s, 1.0, 2.3, 11.3, 3.4, [
    "El recuento del dinero se rinde y se concilia con tesorería.",
    "Descarga segura: tacos de seguridad, cortinas y zona segura.",
    "Sanitizado de la unidad obligatorio todos los días.",
    "Checklist de retorno (CloudFleet) en cada regreso.",
    "Tiempos internos: verificación física y cierre financiero, ≤ 40 min cada uno.",
], size=15.5, gap=11)
box(s, 1.0, 6.05, 6.5, 0.02, fill=CARD)
text(s, 1.0, 6.25, 11.3, 0.7, [[("Un cierre prolijo del reparto asegura el control del dinero, la calidad del producto y la unidad lista para el día siguiente.", 14, AMBER, True)]], line_spacing=1.1)

prs.save("/tmp/material-1.3.pptx")
print(f"OK /tmp/material-1.3.pptx — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
